import io
from pptx import Presentation
from typing import List, Dict, Any, Union
from app.rag.parsers.base import BaseParser


class PptxParser(BaseParser):
    """
    Parser for extracting text from Microsoft PowerPoint (.pptx) slide decks.
    Each slide is treated as a physical page.
    """
    def parse(self, file_input: Union[bytes, str]) -> List[Dict[str, Any]]:
        if isinstance(file_input, str):
            prs = Presentation(file_input)
        else:
            prs = Presentation(io.BytesIO(file_input))
        pages = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())

            pages.append({
                "text": "\n".join(slide_text_parts),
                "page_number": slide_idx + 1
            })

        return pages
